#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
precheck.py — deck 铁律机器预检（转图人眼 QA 之前先跑）

把 CLAUDE.md 铁律里能程序化的部分变成一条命令：
  1) 出画布      元素越出 13.33"×7.5" 画布（pptxgenjs 出界坐标照写不裁剪）
  2) 边距纪律    内容左缘 < 0.5" 或右缘 > 12.8"（装饰性出血件单独提示）
  3) 字号下限    正文 ≥13pt、标签 ≥12pt、页脚带（y>6.85"）豁免 9–10pt
  4) 半角标点    中文字符紧邻半角 , ; : ! ? ( ) 一律报错
  5) 文本互撞    文本框×文本框 / 图片×文本框 重叠（色块垫底属正常版式，不报）

用法:
  python3 scripts/precheck.py build/xxx.pptx            # 全量检查
  python3 scripts/precheck.py build/xxx.pptx --slide 3  # 只查第 3 页

输出按页分组、按级别标记:
  [错误] 必改;  [警告] 人工确认;  [提示] 供参考
退出码: 有 [错误] → 1, 否则 0
"""
import argparse
import re
import sys
from pptx import Presentation
from pptx.util import Emu

PAGE_W, PAGE_H = 13.333, 7.5     # LAYOUT_WIDE(英寸)
MARGIN_L, MARGIN_R = 0.5, 12.8   # 内容带
EPS = 0.06                        # 容差
FOOTER_Y = 6.45                   # 此线以下按页脚/引用带豁免小字号
BODY_MIN, LABEL_MIN, FOOT_MIN = 13, 12, 9

# 中文字符紧邻半角标点(全角规则铁律);允许 . 和 - 出现在数字/英文场景
HALF_PUNCT = re.compile(r"[一-鿿][,;:!?()]|[,;:!?()][一-鿿]")


def inches(v):
    return Emu(v).inches if v is not None else None


def shape_box(sh):
    try:
        return (inches(sh.left), inches(sh.top), inches(sh.width), inches(sh.height))
    except Exception:
        return (None, None, None, None)


def iter_text_runs(sh):
    if not sh.has_text_frame:
        return
    for para in sh.text_frame.paragraphs:
        for run in para.runs:
            yield run


def shape_text(sh):
    if not sh.has_text_frame:
        return ""
    return "\n".join(p.text for p in sh.text_frame.paragraphs)


def overlap_area(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    if None in a or None in b:
        return 0.0
    w = min(ax + aw, bx + bw) - max(ax, bx)
    h = min(ay + ah, by + bh) - max(ay, by)
    return max(0.0, w) * max(0.0, h)


def check_deck(path, only_slide=None):
    prs = Presentation(path)
    n_err = n_warn = 0
    for idx, slide in enumerate(prs.slides, 1):
        if only_slide and idx != only_slide:
            continue
        errs, warns, infos = [], [], []
        shapes = []
        for sh in slide.shapes:
            x, y, w, h = shape_box(sh)
            if x is None:
                continue
            kind = sh.shape_type
            is_line = "LINE" in str(kind) or (h is not None and h < 0.02) or (w is not None and w < 0.02)
            has_text = sh.has_text_frame and shape_text(sh).strip()
            shapes.append(dict(sh=sh, x=x, y=y, w=w, h=h, line=is_line,
                               text=bool(has_text), pic="PICTURE" in str(kind)))

            label = (shape_text(sh).strip().replace("\n", " ")[:18] or str(kind))

            # 1) 出画布
            if x + w > PAGE_W + EPS or y + h > PAGE_H + EPS or x < -EPS or y < -EPS:
                # 完全出界(装饰出血)与文字出界区别对待
                if has_text:
                    errs.append(f"文字元素出画布「{label}」 x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
                else:
                    infos.append(f"装饰元素出血出画布「{label}」(封面圆等属正常) x={x:.2f} y={y:.2f}")

            # 2) 边距(只查带文字的内容件;短小标注如图表刻度豁免)
            if has_text and not is_line:
                txt = shape_text(sh).strip()
                pts = [r.font.size.pt for r in iter_text_runs(sh) if r.font.size]
                tiny_note = len(txt) <= 8 and pts and max(pts) <= 10
                if not tiny_note:
                    if x < MARGIN_L - EPS and x > -EPS:
                        warns.append(f"左缘越过 0.5\" 「{label}」 x={x:.2f}")
                    if x + w > MARGIN_R + EPS and x + w < PAGE_W + EPS:
                        warns.append(f"右缘越过 12.8\" 「{label}」 右={x + w:.2f}")

            # 3) 字号下限
            if has_text:
                in_footer = y >= FOOTER_Y
                for run in iter_text_runs(sh):
                    if run.font.size is None or not run.text.strip():
                        continue
                    pt = run.font.size.pt
                    frag = run.text.strip()[:12]
                    if in_footer:
                        if pt < FOOT_MIN:
                            warns.append(f"页脚字号 {pt:.0f}pt < {FOOT_MIN}pt 「{frag}」")
                    elif pt < LABEL_MIN:
                        if pt >= FOOT_MIN and len(run.text.strip()) <= 8:
                            infos.append(f"小字短标注 {pt:.0f}pt 「{frag}」(图表刻度/数据标注属正常)")
                        elif pt >= FOOT_MIN:
                            warns.append(f"字号 {pt:.0f}pt 低于标签下限 {LABEL_MIN}pt 「{frag}」")
                        else:
                            errs.append(f"字号 {pt:.0f}pt 过小 「{frag}」")
                    elif pt < BODY_MIN:
                        infos.append(f"字号 {pt:.0f}pt(标签档,正文须≥{BODY_MIN}) 「{frag}」")

            # 4) 半角标点
            if has_text:
                for m in HALF_PUNCT.finditer(shape_text(sh)):
                    errs.append(f"半角标点混入中文 「…{m.group(0)}…」")

        # 5) 文本互撞(文本×文本 / 图×文本)
        content = [s for s in shapes if not s["line"]]
        for i in range(len(content)):
            for j in range(i + 1, len(content)):
                a, b = content[i], content[j]
                if not ((a["text"] and b["text"]) or (a["pic"] and b["text"]) or (a["text"] and b["pic"])):
                    continue
                area = overlap_area((a["x"], a["y"], a["w"], a["h"]), (b["x"], b["y"], b["w"], b["h"]))
                small = min(a["w"] * a["h"], b["w"] * b["h"]) or 1
                if area / small > 0.18:
                    la = shape_text(a["sh"]).strip().replace("\n", " ")[:12] or "图/形"
                    lb = shape_text(b["sh"]).strip().replace("\n", " ")[:12] or "图/形"
                    warns.append(f"疑似互撞({area / small:.0%}) 「{la}」×「{lb}」(卡内排版正常层叠请忽略)")

        if errs or warns or infos:
            print(f"\n== Slide {idx} ==")
            for e in errs:
                print(f"  [错误] {e}")
            for w_ in warns:
                print(f"  [警告] {w_}")
            for i_ in infos:
                print(f"  [提示] {i_}")
        n_err += len(errs)
        n_warn += len(warns)

    print(f"\n—— 预检完成: {n_err} 错误 / {n_warn} 警告 ——")
    if n_err == 0 and n_warn == 0:
        print("全绿。接着转图跑人眼 QA(预检不能替代肉眼)。")
    return 1 if n_err else 0


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slide", type=int, default=None)
    args = ap.parse_args()
    sys.exit(check_deck(args.pptx, args.slide))
